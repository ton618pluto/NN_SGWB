from copy import deepcopy

from pptx import Presentation


SRC = "组会汇报_沈益民_v6_压缩背景页.pptx"
OUT = "组会汇报_沈益民_v8_同款样式TimeFlowV2结构页.pptx"


def duplicate_slide(prs, index):
    source = prs.slides[index]
    blank_slide = prs.slides.add_slide(prs.slide_layouts[0])

    for shape in list(blank_slide.shapes):
        el = shape.element
        el.getparent().remove(el)

    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        blank_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return blank_slide


def set_text(shape, text):
    if not shape.has_text_frame:
        return
    shape.text = text


def build_timeflow_slide(prs):
    slide = duplicate_slide(prs, 8)  # duplicate page 9

    replacements = {
        "Model Architecture": "Model Architecture",
        "SimpleCNN1D  →  Embedding  →  GWFlowModel  →  z₀ ~ 𝒩(0, I)": "TimeFlow V2  →  TemporalEncoder1D  →  GWFlowModelV2  →  p(θ | x)",
        "Input": "Input",
        "Conv1d  1→16": "H1 + L1",
        "BN · ReLU · Pool": "[2, 524288]",
        "Conv1d  16→32": "scale ×1e23",
        "Conv1d  32→64": "channel centering",
        "MaxPool  64×512": "dual-channel raw strain",
        "SimpleCNN1D": "Input",
        "Flatten · FC(256)→10": "Stem: Conv1d + GN + GELU",
        "z  (context ∈ ℝ⁵¹²)": "Residual Blocks × 7",
        "Embedding": "TemporalEncoder1D",
        "GWFlowModel  ·  Normalizing Flow × 8": "GWFlowModelV2  ·  Conditional Flow × 6",
        "Flow Layer 1  ·  Permutation + MAF": "Flow Layer  ·  Permutation + Spline AR",
        "Flow Layer 2 – 8  (× 7)": "context_features = 256",
        "z₀ ~ 𝒩(0, I)  ·  Standard Normal Prior": "Base: StandardNormal(10)",
        "Loss:  −log p(θ | x)  (NLL)": "Loss:  −log p(θ | x)  (NLL)",
        "GWFlowModel": "GWFlowModelV2",
        "θ ∈ ℝ¹⁰": "θ ∈ ℝ¹⁰",
        "Conv1d": "Dual Input",
        "Flatten+FC": "Encoder",
        "Context z": "Context z",
        "Flow Layer": "Flow Layer",
        "Prior z₀": "Posterior",
        "Loss": "Loss",
    }

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        current = shape.text.strip()
        if current in replacements:
            set_text(shape, replacements[current])

    # refine a few labels to better fit the old layout
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        current = shape.text.strip()
        if current == "scale ×1e23":
            shape.text_frame.paragraphs[0].font.size = None
        if current == "channel centering":
            shape.text_frame.paragraphs[0].font.size = None

    return slide


def insert_after(prs, slide, position):
    slide_ids = prs.slides._sldIdLst
    new_id = slide_ids[-1]
    slide_ids.remove(new_id)
    slide_ids.insert(position, new_id)


def main():
    prs = Presentation(SRC)
    slide = build_timeflow_slide(prs)
    insert_after(prs, slide, 9)  # after page 9
    prs.save(OUT)


if __name__ == "__main__":
    main()
