from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor


SRC = "组会汇报_沈益民_v4_backup_before_edit.pptx"
OUT = "组会汇报_沈益民_v6_压缩背景页.pptx"


def set_run_format(run, size, bold=False, color=None, name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, align=None, font="Microsoft YaHei"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run_format(r, size=size, bold=bold, color=color, name=font)
    if align is not None:
        p.alignment = align
    return shape


def add_bullets(slide, left, top, width, height, items, font_size, space_after=6):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = f"• {item}"
        set_run_format(r, font_size, False, RGBColor(58, 69, 84))
    return shape


def add_card(slide, left, top, width, height, fill_rgb=RGBColor(255, 255, 255)):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_rgb
    card.line.color.rgb = RGBColor(213, 220, 230)
    card.line.width = Pt(1.0)
    return card


def add_tag(slide, left, top, text, color, width=1.45):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    p = tag.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run_format(r, 13, True, RGBColor(255, 255, 255))
    p.alignment = 1
    return tag


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):
        element = shape.element
        element.getparent().remove(element)

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(247, 249, 252)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.24))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(34, 76, 132)
    bar.line.fill.background()

    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.50), Inches(0.38), Inches(0.58), Inches(0.38))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(34, 76, 132)
    box.line.fill.background()
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "01"
    set_run_format(r, 18, True, RGBColor(255, 255, 255), "Arial")
    p.alignment = 1

    add_textbox(slide, Inches(1.18), Inches(0.34), Inches(6.2), Inches(0.42), "研究背景 — 课题意义及挑战", 22, RGBColor(31, 47, 70), True)
    add_textbox(slide, Inches(1.20), Inches(0.76), Inches(8.8), Inches(0.28), "SGWB 参数估计值得做，但在物理背景、数据规模与训练耗时上都存在现实困难", 10.5, RGBColor(110, 123, 139))

    add_card(slide, Inches(0.62), Inches(1.18), Inches(5.95), Inches(4.82))
    add_card(slide, Inches(6.78), Inches(1.18), Inches(5.95), Inches(4.82))

    add_tag(slide, 0.88, 1.38, "课题意义", RGBColor(34, 76, 132))
    add_tag(slide, 7.05, 1.38, "主要挑战", RGBColor(166, 86, 40))

    significance = [
        "SGWB 仍未被直接探测到，研究它有助于理解宇宙演化、双黑洞种群及早期宇宙信息。",
        "深度学习在 SGWB 参数估计中的应用仍较少，可探索新的高效推断路径。",
        "传统数值求解与贝叶斯推断代价高，而深度学习可把复杂计算前移到训练阶段，在推理时显著提速。",
    ]
    add_bullets(slide, Inches(0.95), Inches(1.86), Inches(5.18), Inches(2.65), significance, 13.6, 7)

    value_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.72), Inches(5.18), Inches(0.92))
    value_band.fill.solid()
    value_band.fill.fore_color.rgb = RGBColor(236, 244, 255)
    value_band.line.color.rgb = RGBColor(180, 204, 235)
    p = value_band.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "核心价值：为复杂的 SGWB 参数反演问题提供“训练耗时、推理快速”的可行方案。"
    set_run_format(r, 13.2, True, RGBColor(38, 78, 132))
    p.alignment = 1

    challenges = [
        "物理背景要求高：需要同时理解恒星、黑洞、引力波及相关种群模型。",
        "相关文献较少：SGWB 与深度学习交叉研究仍然有限。",
        "训练集规模大、特征维度高：数据生成、加载与训练都非常耗时。",
        "目标标签为 10 维超参数，多输出回归进一步提升了建模复杂度。",
    ]
    add_bullets(slide, Inches(7.10), Inches(1.86), Inches(5.00), Inches(2.60), challenges, 12.9, 5)

    stat_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.02), Inches(4.56), Inches(5.15), Inches(0.60))
    stat_band.fill.solid()
    stat_band.fill.fore_color.rgb = RGBColor(252, 245, 232)
    stat_band.line.color.rgb = RGBColor(232, 205, 165)
    p = stat_band.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "例：若仅生成 1 天数据、标签取 10000 组"
    set_run_format(r, 12.8, True, RGBColor(126, 82, 29))
    p.alignment = 1

    stats = [("数据生成", "约 1 个月"), ("数据加载", "约 6 小时"), ("1 个 epoch", "约 10 小时")]
    x_positions = [7.04, 8.78, 10.52]
    widths = [1.5, 1.5, 1.55]
    for (title, value), x, w in zip(stats, x_positions, widths):
        card = add_card(slide, Inches(x), Inches(5.22), Inches(w), Inches(0.95), RGBColor(245, 247, 251))
        tf = card.text_frame
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        set_run_format(r1, 11.4, True, RGBColor(81, 91, 110))
        p1.alignment = 1
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = value
        set_run_format(r2, 13.2, True, RGBColor(120, 78, 148))
        p2.alignment = 1

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.34), Inches(11.22), Inches(0.34))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(239, 241, 247)
    footer.line.color.rgb = RGBColor(210, 214, 223)
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "结论：训练集规模难以简单确定——过大则训练不可行，过小则模型学习有限。"
    set_run_format(r, 11.5, True, RGBColor(90, 97, 112))
    p.alignment = 1

    return slide


def main():
    prs = Presentation(SRC)
    build_slide(prs)
    slide_ids = prs.slides._sldIdLst
    new_id = slide_ids[-1]
    slide_ids.remove(new_id)
    slide_ids.insert(4, new_id)
    prs.save(OUT)


if __name__ == "__main__":
    main()
